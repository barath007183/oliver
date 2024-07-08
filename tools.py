import json
import os

import requests
from bs4 import BeautifulSoup
from crewai import Agent, Task
from langchain.tools import tool
from langchain_community.chat_models import BedrockChat
from pptx import Presentation
from unstructured.partition.html import partition_html

from helpers import CLAUDE_V3_SONNET_MODEL_ID, get_client


class SearchTools:

    @tool("Search internet")
    def search_internet(query):
        """Useful to search the internet about a given topic and return relevant
        results."""
        return SearchTools.search(query)

    def search(query, n_results=5):
        url = "https://google.serper.dev/search"
        payload = json.dumps({"q": query})
        headers = {
            "X-API-KEY": os.environ["SERPER_API_KEY"],
            "content-type": "application/json",
        }
        response = requests.request("POST", url, headers=headers, data=payload)
        results = response.json()["organic"]
        string = []
        for result in results[:n_results]:
            try:
                string.append(
                    "\n".join(
                        [
                            f"Title: {result['title']}",
                            f"Link: {result['link']}",
                            f"Snippet: {result['snippet']}",
                            "\n-----------------",
                        ]
                    )
                )
            except KeyError:
                next

        content = "\n".join(string)
        return f"\nSearch result: {content}\n"


class BrowserTools:

    @tool("Scrape website content")
    def scrape_and_summarize_website(website):
        """Useful to scrape and summarize a website content, just pass a string with
        only the full url, no need for a final slash `/`, eg: https://google.com or https://clearbit.com/about-us
        """
        client = get_client(
            access_key_id=os.getenv("BEDROCK_ACCESS_KEY_ID"),
            secret_access_key=os.getenv("BEDROCK_SECRET_ACCESS_KEY"),
            region=os.getenv("BEDROCK_REGION"),
        )
        llm = BedrockChat(
            client=client,
            model_id=CLAUDE_V3_SONNET_MODEL_ID,
            model_kwargs={
                "max_tokens": 8000,
            },
        )
        response = requests.request("GET", website)
        elements = partition_html(text=response.text)
        content = "\n\n".join([str(el) for el in elements])
        content = [content[i : i + 8000] for i in range(0, len(content), 8000)]
        summaries = []
        for chunk in content:
            agent = Agent(
                role="Principal Researcher",
                goal="Do amazing researches and summaries based on the content you are working with",
                backstory="You're a Principal Researcher at a big company and you need to do a research about a given topic.",
                llm=llm,
                allow_delegation=False,
            )
            task = Task(
                agent=agent,
                description=f"Analyze and make a LONG summary the content bellow, make sure to include the ALL relevant information in the summary, return only the summary nothing else.\n\nCONTENT\n----------\n{chunk}",
                expected_output="A long summary of the content",
            )
            summary = task.execute()
            summaries.append(summary)
            content = "\n\n".join(summaries)
        return f"\nScrapped Content: {content}\n"


@tool("Create a PPT presentation")
def create_ppt(slide_content):
    """
    Useful to create a PPT presentation from slide content.

    Takes a string with the slide content in the following format:
    <slides>
        <slide>
            <type>title_slide</type>
            <title>Title of the slide</title>
        </slide>
        <slide>
            <type>section_header</type>
            <title>Title of the slide</title>
        </slide>
        <slide>
            <type>content_slide</type>
            <title>Title of the slide</title>
            <content>Content of the slide</content>
        </slide>
    </slides>
    """
    root = Presentation()

    """ Ref for slide types:  
    0 ->  title and subtitle 
    1 ->  title and content 
    2 ->  section header 
    3 ->  two content 
    4 ->  Comparison 
    5 ->  Title only  
    6 ->  Blank 
    7 ->  Content with caption 
    8 ->  Pic with caption 
    """
    soup = BeautifulSoup(slide_content, "lxml")
    slide_content = [
        {
            "type": slide.find("type").text,
            "title": slide.find("title").text,
            "content": slide.find("content").text if slide.find("content") else "",
        }
        for slide in soup.find_all("slide")
    ]

    for slide in slide_content:
        if slide["type"] == "title_slide":
            slide_layout = root.slide_layouts[0]
            _slide = root.slides.add_slide(slide_layout)
            _slide.shapes.title.text = slide["title"]
        elif slide["type"] == "section_header":
            slide_layout = root.slide_layouts[2]
            _slide = root.slides.add_slide(slide_layout)
            _slide.shapes.title.text = slide["title"]
        elif slide["type"] == "content_slide":
            slide_layout = root.slide_layouts[1]
            _slide = root.slides.add_slide(slide_layout)
            _slide.shapes.title.text = slide["title"]
            _slide.placeholders[1].text = slide["content"]

    root.save("presentation.pptx")
    return "Presentation created successfully!"
